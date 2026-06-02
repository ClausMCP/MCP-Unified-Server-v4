#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
MCP Office Reader v3.7 – Secure extraction for Word, Excel, PowerPoint.
Enhanced with streaming, Markdown export, comments extraction, DOCX-to-Markdown, and indexed search.
"""
import os
import sys
import time
import pickle
import subprocess
from pathlib import Path
from typing import List, Dict, Optional, Any, Generator
from mcp_shared import (
    _log, normalize_path, _ensure_allowed,
    BaseMCPServer, conversation_memory, dialog_ctx
)
from mcp_rate_limiter import safe_call

# Оригинальные функции с защитой
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

server = BaseMCPServer("office-reader", "3.7")

# Хранилище для индекса (in-memory)
_office_index = {}

# ====================== СУЩЕСТВУЮЩИЕ ФУНКЦИИ ======================

def read_docx(file_path: str, include_headers: bool = True) -> Dict:
    try:
        from docx import Document
    except ImportError:
        return {"error": "python-docx not installed. Install with: pip install python-docx"}
    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "read_docx")
    if not p.is_file():
        return {"error": f"File not found: {file_path}"}
    doc = Document(str(p))
    content = []
    for para in doc.paragraphs:
        content.append({
            "text": para.text,
            "style": para.style.name if include_headers else None,
            "alignment": para.alignment.name if para.alignment else None
        })
    tables = []
    for i, table in enumerate(doc.tables):
        table_data = [[cell.text for cell in row.cells] for row in table.rows]
        tables.append({"table_index": i, "data": table_data})
    result = {"path": str(p), "filename": p.name, "paragraphs": content, "tables": tables}
    conversation_memory.add(
        op="read_docx", paths={"file": str(p)}, status="success", dialog=dialog_ctx.get(),
        context=f"Extracted {len(content)} paragraphs and {len(tables)} tables from {p.name}"
    )
    return result

def read_excel(file_path: str, sheet_name: Optional[str] = None, max_rows: int = 1000) -> Dict:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return {"error": "openpyxl not installed. Install with: pip install openpyxl"}
    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "read_excel")
    if not p.is_file():
        return {"error": f"File not found: {file_path}"}
    wb = load_workbook(str(p), read_only=True, data_only=True)
    ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
    if ws is None:
        return {"error": f"Sheet '{sheet_name}' not found"}
    data = []
    headers = None
    row_count = 0
    for row in ws.iter_rows(values_only=True):
        row_count += 1
        if row_count > max_rows:
            break
        if headers is None:
            headers = list(row)
            continue
        row_dict = {h: v for h, v in zip(headers, row) if h is not None}
        data.append(row_dict)
    wb.close()
    result = {"path": str(p), "sheet": ws.title, "headers": headers, "rows_extracted": len(data), "data": data}
    conversation_memory.add(
        op="read_excel", paths={"file": str(p)}, status="success", dialog=dialog_ctx.get(),
        context=f"Extracted {len(data)} rows from sheet '{ws.title}'"
    )
    return result

def read_pptx(file_path: str, slide_numbers: Optional[List[int]] = None) -> Dict:
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed. Install with: pip install python-pptx"}
    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "read_pptx")
    if not p.is_file():
        return {"error": f"File not found: {file_path}"}
    prs = Presentation(str(p))
    slides_data = []
    for i, slide in enumerate(prs.slides):
        slide_idx = i + 1
        if slide_numbers and slide_idx not in slide_numbers:
            continue
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(paragraph.text)
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        except Exception:
            pass
        slides_data.append({
            "slide_number": slide_idx,
            "text": "\n".join(texts),
            "notes": notes
        })
    result = {"path": str(p), "total_slides": len(prs.slides), "extracted_slides": slides_data}
    conversation_memory.add(
        op="read_pptx", paths={"file": str(p)}, status="success", dialog=dialog_ctx.get(),
        context=f"Extracted text from {len(slides_data)} slides in {p.name}"
    )
    return result

def export_to_pdf(input_path: str, output_path: str) -> Dict:
    src = Path(normalize_path(input_path))
    dst = Path(normalize_path(output_path))
    _ensure_allowed(src, "export_to_pdf")
    _ensure_allowed(dst.parent, "export_to_pdf")
    if not src.is_file():
        return {"error": f"Input file not found: {input_path}"}
    cmd = ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", str(dst.parent), str(src)]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=90)
        if result.returncode != 0:
            alt_cmd = ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(dst.parent), str(src)]
            result = subprocess.run(alt_cmd, capture_output=True, text=True, timeout=90)
        if result.returncode != 0:
            return {"error": f"LibreOffice conversion failed: {result.stderr.strip()}"}
        expected_out = dst.parent / f"{src.stem}.pdf"
        if expected_out.exists():
            if expected_out != dst:
                expected_out.rename(dst)
        conversation_memory.add(
            op="export_to_pdf", paths={"src": str(src), "dst": str(dst)},
            status="success", dialog=dialog_ctx.get(), context=f"Converted {src.name} to PDF"
        )
        return {"status": "success", "output": str(dst)}
        return {"error": "Conversion completed but output file not found."}
    except FileNotFoundError:
        return {"error": "LibreOffice/soffice not found. Install it to enable PDF export."}
    except Exception as e:
        return {"error": str(e)}

# ====================== НОВЫЕ ФУНКЦИИ (v3.6) ======================

def extract_text_from_docx_safe(file_path: str, max_size_mb: int = 50) -> Dict:
    """Безопасное извлечение текста из DOCX с ограничением размера."""
    def _internal():
        p = Path(normalize_path(file_path))
        _ensure_allowed(p, "extract_docx")
        if not p.exists():
            return {"error": "File not found"}
        if p.stat().st_size > max_size_mb * 1024 * 1024:
            return {"error": f"File exceeds {max_size_mb} MB"}
        if not DOCX_AVAILABLE:
            return {"error": "python-docx not installed"}
        doc = Document(str(p))
        text = "\n".join(para.text for para in doc.paragraphs)
        tables = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append([cell.text for cell in row.cells])
            tables.append(rows)
        return {"text": text[:50000], "tables": tables[:10], "total_chars": len(text), "path": str(p)}
    return safe_call("office_docx", _internal)

def extract_excel_sheet_safe(file_path: str, sheet_name: Optional[str] = None, max_rows: int = 1000) -> Dict:
    def _internal():
        p = Path(normalize_path(file_path))
        _ensure_allowed(p, "extract_excel")
        if not XLSX_AVAILABLE:
            return {"error": "openpyxl not installed"}
        wb = load_workbook(str(p), read_only=True, data_only=True)
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
        if ws is None:
            return {"error": "Sheet not found"}
        data = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows:
                break
            data.append([str(cell) if cell is not None else "" for cell in row])
        wb.close()
        return {"sheet": ws.title, "data": data, "rows": len(data), "path": str(p)}
    return safe_call("office_excel", _internal)

def search_in_office_files(folder_path: str, keyword: str, extensions: List[str] = None) -> Dict:
    def _internal():
        root = Path(normalize_path(folder_path))
        _ensure_allowed(root, "search_office")
        if not root.is_dir():
            return {"error": "Not a directory"}
        exts = extensions or [".docx", ".xlsx", ".pptx"]
        results = []
        for ext in exts:
            for f in root.rglob(f"*{ext}"):
                try:
                    if ext == ".docx" and DOCX_AVAILABLE:
                        doc = Document(str(f))
                        text = " ".join(p.text for p in doc.paragraphs)
                        if keyword.lower() in text.lower():
                            results.append({"file": str(f), "type": "docx"})
                    elif ext == ".xlsx" and XLSX_AVAILABLE:
                        wb = load_workbook(str(f), read_only=True)
                        found = False
                        for sheet in wb.worksheets:
                            for row in sheet.iter_rows(values_only=True):
                                if any(keyword.lower() in str(cell).lower() for cell in row if cell):
                                    found = True
                                    break
                            if found:
                                break
                        wb.close()
                        if found:
                            results.append({"file": str(f), "type": "xlsx"})
                    elif ext == ".pptx" and PPTX_AVAILABLE:
                        prs = Presentation(str(f))
                        text = " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
                        if keyword.lower() in text.lower():
                            results.append({"file": str(f), "type": "pptx"})
                except Exception as e:
                    _log(f"Search error in {f}: {e}")
        return {"keyword": keyword, "matches": results, "count": len(results)}
    return safe_call("office_search", _internal)

# ====================== НОВЫЕ ФУНКЦИИ (v3.7) ======================

# 1.1. Потоковое чтение больших Excel (снижение памяти)
def read_excel_streaming(file_path: str, sheet_name: Optional[str] = None, chunk_size: int = 10000) -> Dict:
    """Читает Excel построчно, возвращая генератор. Не держит весь файл в памяти."""
    def _internal():
        from openpyxl import load_workbook
        p = Path(normalize_path(file_path))
        _ensure_allowed(p, "read_excel_streaming")
        if not p.is_file():
            return {"error": f"File not found: {file_path}"}
        
        wb = load_workbook(str(p), read_only=True, data_only=True)
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
        headers = None
        
        def _generator():
            nonlocal headers
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if headers is None:
                    headers = [str(cell) if cell is not None else f"col_{j}" for j, cell in enumerate(row)]
                    continue
                yield dict(zip(headers, row))
                if i % chunk_size == 0:
                    time.sleep(0)  # дать шанс другим корутинам
            wb.close()
        
        return {"generator": _generator(), "sheet": ws.title, "path": str(p)}
    return safe_call("office_excel_stream", _internal)

# 1.2. Извлечение таблиц в Markdown
def excel_to_markdown(file_path: str, sheet_name: Optional[str] = None, max_rows: int = 100) -> Dict:
    """Возвращает Markdown-таблицу из Excel."""
    def _internal():
        p = Path(normalize_path(file_path))
        _ensure_allowed(p, "excel_to_markdown")
        if not p.is_file():
            return {"error": f"File not found: {file_path}"}
        if not PANDAS_AVAILABLE:
            return {"error": "pandas not installed. Install with: pip install pandas"}
        
        import pandas as pd
        df = pd.read_excel(str(p), sheet_name=sheet_name, nrows=max_rows)
        markdown = df.to_markdown(index=False)
        return {"markdown": markdown, "rows": len(df), "path": str(p)}
    return safe_call("office_excel_md", _internal)

# 1.3. Извлечение комментариев и скрытых листов
def read_excel_comments(file_path: str, sheet_name: Optional[str] = None) -> Dict:
    """Возвращает все комментарии (примечания) ячеек."""
    def _internal():
        from openpyxl import load_workbook
        p = Path(normalize_path(file_path))
        _ensure_allowed(p, "read_excel_comments")
        if not p.is_file():
            return {"error": f"File not found: {file_path}"}
        
        wb = load_workbook(str(p), data_only=False)
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
        comments = {}
        for cell in ws._cells.values():
            if cell.comment:
                comments[cell.coordinate] = cell.comment.text
        wb.close()
        return {"sheet": ws.title, "comments": comments, "path": str(p)}
    return safe_call("office_excel_comments", _internal)

# 1.4. Конвертация DOCX в Markdown с сохранением структуры
def docx_to_markdown(file_path: str) -> Dict:
    """Преобразует DOCX в Markdown (заголовки, списки, таблицы)."""
    def _internal():
        from docx import Document
        p = Path(normalize_path(file_path))
        _ensure_allowed(p, "docx_to_markdown")
        if not p.is_file():
            return {"error": f"File not found: {file_path}"}
        
        doc = Document(str(p))
        md_lines = []
        for para in doc.paragraphs:
            style = para.style.name
            text = para.text
            if style.startswith('Heading'):
                level = style[-1] if style[-1].isdigit() else '1'
                md_lines.append(f"{'#' * int(level)} {text}")
            elif style == 'List Bullet':
                md_lines.append(f"- {text}")
            elif style == 'List Number':
                md_lines.append(f"1. {text}")
            else:
                md_lines.append(text)
        
        # Таблицы в Markdown
        for table in doc.tables:
            md_lines.append("")
            # Заголовки таблицы
            headers = [cell.text for cell in table.rows[0].cells]
            md_lines.append("| " + " | ".join(headers) + " |")
            md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
            # Данные
            for row in table.rows[1:]:
                cells = [cell.text for cell in row.cells]
                md_lines.append("| " + " | ".join(cells) + " |")
            md_lines.append("")
        
        markdown = "\n".join(md_lines)
        return {"markdown": markdown, "paragraphs": len(doc.paragraphs), "tables": len(doc.tables), "path": str(p)}
    return safe_call("office_docx_md", _internal)

# 1.5. Улучшенный поиск по содержимому с использованием индекса
def build_office_index(folder_path: str, extensions: List[str] = None) -> Dict:
    """Строит инвертированный индекс для быстрого поиска в офисных файлах."""
    global _office_index
    def _internal():
        root = Path(normalize_path(folder_path))
        _ensure_allowed(root, "build_index")
        if not root.is_dir():
            return {"error": "Not a directory"}
        
        exts = extensions or [".docx", ".xlsx", ".pptx"]
        index = {}
        file_count = 0
        
        for ext in exts:
            for f in root.rglob(f"*{ext}"):
                try:
                    file_count += 1
                    file_path = str(f)
                    words = set()
                    
                    if ext == ".docx" and DOCX_AVAILABLE:
                        doc = Document(file_path)
                        text = " ".join(p.text for p in doc.paragraphs)
                        words = set(text.lower().split())
                    elif ext == ".xlsx" and XLSX_AVAILABLE:
                        wb = load_workbook(file_path, read_only=True)
                        all_text = []
                        for sheet in wb.worksheets:
                            for row in sheet.iter_rows(values_only=True):
                                all_text.extend([str(cell).lower() for cell in row if cell])
                        wb.close()
                        words = set(" ".join(all_text).split())
                    elif ext == ".pptx" and PPTX_AVAILABLE:
                        prs = Presentation(file_path)
                        text = " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
                        words = set(text.lower().split())
                    
                    # Инвертированный индекс: слово -> список файлов
                    for word in words:
                        if word not in index:
                            index[word] = []
                        index[word].append(file_path)
                except Exception as e:
                    _log(f"Index error in {f}: {e}")
        
        _office_index = index
        return {"status": "success", "files_indexed": file_count, "unique_words": len(index), "path": str(root)}
    return safe_call("office_build_index", _internal)

def search_with_index(folder_path: str, keyword: str, extensions: List[str] = None) -> Dict:
    """Быстрый поиск с использованием предварительно построенного индекса."""
    global _office_index
    def _internal():
        if not _office_index:
            return {"error": "Index not built. Run build_office_index first."}
        
        keyword_lower = keyword.lower()
        matches = []
        
        # Поиск точного совпадения
        if keyword_lower in _office_index:
            matches = _office_index[keyword_lower]
        
        # Поиск частичного совпадения
        partial_matches = []
        for word, files in _office_index.items():
            if keyword_lower in word:
                partial_matches.extend(files)
        
        # Уникальные результаты
        all_matches = list(set(matches + partial_matches))
        
        return {"keyword": keyword, "matches": all_matches, "count": len(all_matches), "search_type": "indexed"}
    return safe_call("office_search_index", _internal)

# ====================== РЕГИСТРАЦИЯ ИНСТРУМЕНТОВ ======================

# Новые инструменты v3.7
server.register_tool("read_excel_streaming", {
    "description": "Потоковое чтение больших Excel файлов (снижение памяти)",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "sheet_name": {"type": "string"},
            "chunk_size": {"type": "integer", "default": 10000}
        },
        "required": ["file_path"]
    }
}, lambda **kw: read_excel_streaming(kw["file_path"], kw.get("sheet_name"), kw.get("chunk_size", 10000)))

server.register_tool("excel_to_markdown", {
    "description": "Извлечение Excel-таблиц в Markdown-формат",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "sheet_name": {"type": "string"},
            "max_rows": {"type": "integer", "default": 100}
        },
        "required": ["file_path"]
    }
}, lambda **kw: excel_to_markdown(kw["file_path"], kw.get("sheet_name"), kw.get("max_rows", 100)))

server.register_tool("read_excel_comments", {
    "description": "Извлечение комментариев из ячеек Excel",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "sheet_name": {"type": "string"}
        },
        "required": ["file_path"]
    }
}, lambda **kw: read_excel_comments(kw["file_path"], kw.get("sheet_name")))

server.register_tool("docx_to_markdown", {
    "description": "Конвертация DOCX в Markdown с сохранением структуры",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"}
        },
        "required": ["file_path"]
    }
}, lambda **kw: docx_to_markdown(kw["file_path"]))

server.register_tool("build_office_index", {
    "description": "Построение индекса для быстрого поиска в офисных файлах",
    "inputSchema": {
        "type": "object",
        "properties": {
            "folder_path": {"type": "string"},
            "extensions": {"type": "array", "items": {"type": "string"}}
        },
        "required": ["folder_path"]
    }
}, lambda **kw: build_office_index(kw["folder_path"], kw.get("extensions")))

server.register_tool("search_with_index", {
    "description": "Быстрый поиск с использованием индекса (требует build_office_index)",
    "inputSchema": {
        "type": "object",
        "properties": {
            "folder_path": {"type": "string"},
            "keyword": {"type": "string"}
        },
        "required": ["folder_path", "keyword"]
    }
}, lambda **kw: search_with_index(kw["folder_path"], kw["keyword"]))

# Старые инструменты для совместимости
server.register_tool("extract_docx", {
    "description": "Извлечь текст и таблицы из DOCX с ограничением размера",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "max_size_mb": {"type": "integer", "default": 50}
        },
        "required": ["file_path"]
    }
}, lambda **kw: extract_text_from_docx_safe(kw["file_path"], kw.get("max_size_mb", 50)))

server.register_tool("extract_excel", {
    "description": "Извлечь данные из Excel (поддержка больших файлов)",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "sheet_name": {"type": "string"},
            "max_rows": {"type": "integer", "default": 1000}
        },
        "required": ["file_path"]
    }
}, lambda **kw: extract_excel_sheet_safe(kw["file_path"], kw.get("sheet_name"), kw.get("max_rows", 1000)))

server.register_tool("search_office_files", {
    "description": "Поиск ключевого слова во всех офисных файлах в папке",
    "inputSchema": {
        "type": "object",
        "properties": {
            "folder_path": {"type": "string"},
            "keyword": {"type": "string"}
        },
        "required": ["folder_path", "keyword"]
    }
}, lambda **kw: search_in_office_files(kw["folder_path"], kw["keyword"]))

server.register_tool("read_docx", {
    "description": "Extract text, tables, and styles from DOCX files",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "include_headers": {"type": "boolean", "default": True}
        },
        "required": ["file_path"]
    }
}, lambda **kw: read_docx(kw["file_path"], kw.get("include_headers", True)))

server.register_tool("read_excel", {
    "description": "Extract sheet data from Excel files to JSON structure",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "sheet_name": {"type": "string"},
            "max_rows": {"type": "integer", "default": 1000}
        },
        "required": ["file_path"]
    }
}, lambda **kw: read_excel(kw["file_path"], kw.get("sheet_name"), kw.get("max_rows", 1000)))

server.register_tool("read_pptx", {
    "description": "Extract text and speaker notes from PowerPoint presentations",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "slide_numbers": {"type": "array", "items": {"type": "integer"}}
        },
        "required": ["file_path"]
    }
}, lambda **kw: read_pptx(kw["file_path"], kw.get("slide_numbers")))

server.register_tool("export_to_pdf", {
    "description": "Convert supported office documents to PDF via LibreOffice",
    "inputSchema": {
        "type": "object",
        "properties": {
            "input_path": {"type": "string"},
            "output_path": {"type": "string"}
        },
        "required": ["input_path", "output_path"]
    }
}, lambda **kw: export_to_pdf(kw["input_path"], kw["output_path"]))

if __name__ == "__main__":
    server.run()