#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
MCP Filesystem Search v4.2 (Context-Isolated & Optimized)
Removes asyncio dependency. Uses ThreadPoolExecutor for non-blocking IO.
Compatible with mcp_shared v4.0+ (SQLite). Integrates dialog context isolation.
FIXES: 
 - Removed forced chunk_cache.delete (relies on TTL).
 - analyze_directory rewritten iteratively to prevent RecursionError.
 - get_file_tree rewritten iteratively to prevent RecursionError.
"""
import os
import sys
import json
import re
import time
import hashlib
import fnmatch
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError
from pathlib import Path
from datetime import datetime
from collections import defaultdict
from typing import List, Dict, Optional
from mcp_shared import (
    _log, _format_size, normalize_path, _ensure_allowed,
    CHUNK_SIZE, ENABLE_PAGINATION,
    MAX_READ_BYTES, BaseMCPServer, chunk_cache, conversation_memory,
    list_directory_sync, dialog_ctx, send_progress, is_verbose
)
# Переопределяем таймауты локально с поддержкой env (для долгих операций LM Studio)
SEARCH_TIMEOUT = int(os.environ.get("MCP_SEARCH_TIMEOUT", "3600"))       # было 60
ANALYSIS_TIMEOUT = int(os.environ.get("MCP_ANALYSIS_TIMEOUT", "3600"))   # было 120
_executor = ThreadPoolExecutor(max_workers=8, thread_name_prefix="mcp_search")

# ─── Categorization (static, fast) ───────────────────────────────────────────
def _categorize(path: str) -> str:
    pl = path.lower()
    if any(k in pl for k in ("tools", "utils", "soft", "setup", "install")): return "tools"
    if any(k in pl for k in ("tv", "video", "movie", "movies")): return "media_tv"
    if any(k in pl for k in ("music", "audio", "sound", "mp3", "lossless")): return "media_audio"
    if any(k in pl for k in ("game", "games", "gaming")): return "games"
    if any(k in pl for k in ("doc", "office", "docs", "work", "project")): return "docs"
    if any(k in pl for k in ("backup", "archive", "old", "trash", "temp")): return "archive"
    if any(k in pl for k in ("image", "photo", "picture", "img", "screenshot")): return "images"
    if any(k in pl for k in ("download", "torrent")): return "downloads"
    return "other"

# ─── Sync Walkers ────────────────────────────────────────────────────────────
def _walk_sync(base: str, raw_patterns: list, max_files: int, timeout: int):
    collected = []
    seen = set()
    start = time.time()
    for root, dirs, files in os.walk(base):
        if time.time() - start > timeout:
            break
        dirs[:] = [d for d in dirs if not d.startswith('.') and os.access(os.path.join(root, d), os.R_OK | os.X_OK)]
        for f in files:
            if any(fnmatch.fnmatch(f, pat) for pat in raw_patterns):
                full = os.path.join(root, f)
                if full in seen: continue
                seen.add(full)
                try:
                    st = os.stat(full)
                    collected.append({"name": f, "path": full, "size": st.st_size, "size_human": _format_size(st.st_size), "modified": datetime.fromtimestamp(st.st_mtime).isoformat()})
                except: pass
            if len(collected) >= max_files: return collected, seen
    return collected, seen

def _scandir_sync(base: str, raw_patterns: list):
    collected = []
    seen = set()
    try:
        for entry in os.scandir(base):
            if entry.is_file() and any(fnmatch.fnmatch(entry.name, pat) for pat in raw_patterns):
                full = entry.path
                if full in seen: continue
                seen.add(full)
                try:
                    st = entry.stat(follow_symlinks=False)
                    collected.append({"name": entry.name, "path": full, "size": st.st_size, "size_human": _format_size(st.st_size), "modified": datetime.fromtimestamp(st.st_mtime).isoformat()})
                except: pass
    except PermissionError: pass
    return collected, seen

# ─── Core Search ─────────────────────────────────────────────────────────────
def search_files(path: str, pattern: str, recursive: bool = True, max_files: int = 10000, chunk_id: str = None):
    if chunk_id:
        cached = chunk_cache.get(chunk_id)
        if cached is not None:
            items = cached
            total = len(items)
            start, end = 0, CHUNK_SIZE
            chunk_items = items[start:end]
            has_more = end < total
            resp = {"pattern": pattern, "path": path, "items": chunk_items, "count": len(chunk_items), "total": total}
            if has_more: resp["next_chunk"] = {"chunk_id": chunk_id, "offset": end}
            # Убрано: else: chunk_cache.delete(chunk_id) — очистка теперь только по TTL
            return resp

    p = Path(normalize_path(path))
    try: _ensure_allowed(p, "search_files")
    except PermissionError as e: return {"error": str(e)}

    raw_patterns = [x.strip() for x in pattern.split(';') if x.strip()]
    if not raw_patterns: return {"pattern": pattern, "path": str(p), "items": [], "count": 0, "total": 0}

    base = str(p.resolve())
    try:
        if recursive: future = _executor.submit(_walk_sync, base, raw_patterns, max_files, SEARCH_TIMEOUT)
        else: future = _executor.submit(_scandir_sync, base, raw_patterns)
        collected, seen = future.result(timeout=SEARCH_TIMEOUT)
    except FuturesTimeoutError:
        _log(f"Search timeout after {SEARCH_TIMEOUT}s for {path}")
        collected = []
    except Exception as e:
        _log(f"Search error: {e}")
        collected = []

    if ENABLE_PAGINATION and len(collected) > CHUNK_SIZE:
        cid = f"sch_{int(time.time() * 1000) & 0xFFFFFFFF:08x}"
        chunk_cache.set(cid, collected)
        chunk = collected[:CHUNK_SIZE]
        return {"pattern": pattern, "path": str(p), "items": chunk, "count": len(chunk), "total": len(collected), "chunk_id": cid, "next_chunk": {"chunk_id": cid, "offset": CHUNK_SIZE}}
    return {"pattern": pattern, "path": str(p), "items": collected, "count": len(collected), "total": len(collected)}

def search_content(path: str, query: str, extensions=None, case_sensitive: bool = False, max_files: int = 500):
    if not query or not query.strip(): return [{"error": "Empty search query"}]
    if extensions is None: extensions = (".txt", ".log", ".json", ".py", ".md", ".csv", ".xml", ".ini", ".cfg")
    else: extensions = tuple(e.lower() for e in extensions)

    p = Path(normalize_path(path))
    try: _ensure_allowed(p, "search_content")
    except PermissionError as e: return [{"error": str(e)}]

    flags = 0 if case_sensitive else re.IGNORECASE
    try: pat = re.compile(re.escape(query), flags)
    except re.error: return [{"error": "Invalid regex pattern"}]

    results = []
    start_time = time.time()
    def _content_worker():
        res = []
        for root, dirs, files in os.walk(str(p.resolve())):
            if time.time() - start_time > SEARCH_TIMEOUT: break
            dirs[:] = [d for d in dirs if os.access(os.path.join(root, d), os.R_OK)]
            for f in files:
                if not f.lower().endswith(extensions): continue
                full = os.path.join(root, f)
                try:
                    if os.path.getsize(full) > MAX_READ_BYTES: continue
                    matches = []
                    with open(full, 'r', encoding='utf-8', errors='ignore') as fh:
                        for i, line in enumerate(fh):
                            if pat.search(line):
                                matches.append({"line": i+1, "text": line.strip()[:200]})
                                if len(matches) >= 5: break
                    if matches: res.append({"path": full, "matches": matches, "match_count": len(matches)})
                except: pass
                if len(res) >= max_files: return res
        return res
    try:
        future = _executor.submit(_content_worker)
        results = future.result(timeout=SEARCH_TIMEOUT)
    except FuturesTimeoutError: _log("Content search timeout")
    except Exception as e: _log(f"Content search error: {e}")
    return results

def search_content_advanced(path: str, query: str, use_regex: bool = False, case_sensitive: bool = False,
                            context_lines: int = 0, max_files: int = 100, extensions: List[str] = None,
                            verbose: bool = False, _task_progress=None, **kwargs) -> Dict:
    """
    Advanced search with regex support and context lines.
    Supports progress reporting via send_progress (if verbose=True) or via _task_progress (async tasks).
    """
    root = Path(normalize_path(path))
    try:
        _ensure_allowed(root, "search_content_advanced")
    except PermissionError as e:
        return {"status": "error", "message": str(e)}
    if not root.exists() or not root.is_dir():
        return {"status": "error", "message": "Invalid directory"}

    flags = 0 if case_sensitive else re.IGNORECASE
    try:
        pattern = re.compile(query, flags) if use_regex else re.compile(re.escape(query), flags)
    except re.error as e:
        return {"status": "error", "message": f"Invalid regex: {e}"}

    start_time = time.time()
    results, files_scanned = [], 0
    dialog_id = dialog_ctx.get()

    total_files_estimate = 0
    for _ in root.rglob('*'):
        total_files_estimate += 1
        if total_files_estimate > max_files * 2:
            break
    total_files = min(total_files_estimate, max_files * 2) or max_files

    last_progress_time = 0
    last_progress_processed = 0
    def report_progress(processed, stage="scanning"):
        nonlocal last_progress_time, last_progress_processed
        if _task_progress:
            _task_progress.update_progress(stage, current=processed, total=total_files)
        elif verbose:
            now = time.time()
            if (now - last_progress_time > 2) or (processed - last_progress_processed > 100):
                send_progress(dialog_id, f"Поиск: обработано {processed} из ~{total_files} файлов...")
                last_progress_time = now
                last_progress_processed = processed

    processed = 0
    for entry in root.rglob('*'):
        if processed >= max_files * 2:
            break
        if _task_progress and _task_progress.should_stop():
            _task_progress.save_resume_data({"last_path": str(entry)})
            return None
        if not entry.is_file():
            continue
        if extensions and entry.suffix.lower() not in [e.lower() for e in extensions]:
            processed += 1
            continue
        try:
            if entry.stat().st_size > MAX_READ_BYTES:
                processed += 1
                continue
        except:
            processed += 1
            continue

        processed += 1
        files_scanned += 1
        if files_scanned > max_files:
            break
        if processed % 50 == 0:
            report_progress(processed)

        try:
            with open(entry, 'r', encoding='utf-8', errors='ignore') as f:
                lines = f.readlines()
                for i, line in enumerate(lines):
                    if pattern.search(line):
                        s, e_idx = max(0, i - context_lines), min(len(lines), i + context_lines + 1)
                        results.append({
                            "file": str(entry),
                            "matches": [{
                                "line_number": i+1,
                                "match": line.strip(),
                                "context": [l.strip() for l in lines[s:e_idx]]
                            }]
                        })
                        break
        except Exception:
            continue

        if len(results) >= max_files:
            break

    report_progress(processed, stage="completed")
    result = {
        "status": "success",
        "path": str(root),
        "query": query,
        "use_regex": use_regex,
        "files_scanned": files_scanned,
        "total_matches": len(results),
        "elapsed_sec": round(time.time() - start_time, 2),
        "results": results[:50]
    }
    if _task_progress:
        _task_progress.add_chunk(0, result)
        return None

    conversation_memory.add(
        op="search_content_advanced",
        paths={"path": str(root)},
        status="completed",
        dialog=dialog_id,
        context=f"Searched '{query}' in {files_scanned} files, found {len(results)} matches"
    )
    return result

# ─── NEW: Search by Keywords (AND logic) ─────────────────────────────────────
def search_by_keywords(path: str, keywords: List[str], extensions: List[str] = None,
                       case_sensitive: bool = False, max_files: int = 50) -> Dict:
    if not keywords:
        return {"status": "error", "message": "No keywords provided"}
    pattern_parts = [f"(?=.*{re.escape(kw)})" for kw in keywords]
    regex_pattern = "".join(pattern_parts)
    return search_content_advanced(
        path=path, query=regex_pattern, use_regex=True, case_sensitive=case_sensitive,
        context_lines=0, max_files=max_files, extensions=extensions
    )

# ─── Find Duplicates ─────────────────────────────────────────────────────────
def _file_hash(fpath: str, full: bool = True) -> str:
    h = hashlib.md5()
    try:
        if full:
            with open(fpath, 'rb') as f:
                for chunk in iter(lambda: f.read(65536), b''): h.update(chunk)
        else:
            size = os.path.getsize(fpath)
            with open(fpath, 'rb') as f:
                h.update(f.read(65536))
                if size > 131072: f.seek(size // 2); h.update(f.read(65536))
                if size > 65536: f.seek(max(0, size - 65536)); h.update(f.read(65536))
    except: pass
    return h.hexdigest()

def find_duplicates(path: str, by="hash", min_size_kb=1, max_size_mb=500, extensions=None):
    p = Path(normalize_path(path))
    try: _ensure_allowed(p, "find_duplicates")
    except PermissionError as e: return {"error": str(e)}

    start_time = time.time()
    min_bytes, max_bytes = min_size_kb * 1024, max_size_mb * 1024 * 1024
    files = []
    for root, dirs, filenames in os.walk(str(p.resolve())):
        if time.time() - start_time > ANALYSIS_TIMEOUT: break
        dirs[:] = [d for d in dirs if os.access(os.path.join(root, d), os.R_OK)]
        for f in filenames:
            if extensions and not any(f.lower().endswith(e.lower()) for e in extensions): continue
            full = os.path.join(root, f)
            try:
                size = os.path.getsize(full)
                if min_bytes <= size <= max_bytes: files.append((full, size))
            except: continue

    if by == "size":
        groups = defaultdict(list)
        for f, size in files: groups[size].append(f)
        return [sorted(g) for g in groups.values() if len(g) > 1]

    size_groups = defaultdict(list)
    for f, size in files: size_groups[size].append(f)
    hash_groups = defaultdict(list)
    for size, group in size_groups.items():
        if len(group) == 1: continue
        for f in group:
            try: hash_groups[_file_hash(f, full=(size <= 100 * 1024 * 1024))].append(f)
            except: continue
    return [sorted(g) for g in hash_groups.values() if len(g) > 1]

# ─── File Tree (ITERATIVE, стек, предотвращает RecursionError) ──────────────
def get_file_tree(path: str, max_depth=5, include_files=False):
    """
    Итеративная версия построения дерева каталогов с использованием стека.
    Исключает рекурсию, безопасна для глубоких деревьев (глубина > 1000).
    """
    p = Path(normalize_path(path))
    try:
        _ensure_allowed(p, "get_file_tree")
    except PermissionError as e:
        return [f"⚠️ Access denied: {e}"]

    lines = []
    start_time = time.time()
    lines.append(f"📁 {p}")
    # Стек: (путь, префикс, глубина, флаг_последний_в_родителе)
    stack = [(p, "", 1, True)]

    while stack and (time.time() - start_time) <= ANALYSIS_TIMEOUT:
        cur_path, prefix, depth, is_last = stack.pop()
        if depth > max_depth:
            continue

        try:
            entries = list(os.scandir(cur_path))
        except (PermissionError, OSError):
            lines.append(f"{prefix}⚠️ Access denied")
            continue

        # Сортировка: сначала директории, затем по имени (как в оригинале)
        entries.sort(key=lambda e: (not e.is_dir(follow_symlinks=False), e.name.lower()))
        # Для сохранения порядка (слева направо) помещаем детей в стек в обратном порядке
        children = []
        for i, entry in enumerate(entries):
            is_last_child = (i == len(entries) - 1)
            children.append((entry, is_last_child))
        for entry, is_last_child in reversed(children):
            connector = "└── " if is_last_child else "├── "
            if entry.is_dir(follow_symlinks=False):
                lines.append(f"{prefix}{connector}{entry.name}/")
                child_prefix = prefix + ("    " if is_last_child else "│   ")
                stack.append((Path(entry.path), child_prefix, depth + 1, is_last_child))
            elif include_files:
                lines.append(f"{prefix}{connector}{entry.name}")

    return lines

# ─── Analyze Directory (Iterative, Stack-based) ──────────────────────────────
def analyze_directory(path: str, max_depth=3, group_by="extension"):
    p = Path(normalize_path(path))
    try: _ensure_allowed(p, "analyze_directory")
    except PermissionError as e: return {"error": str(e)}
    if not p.is_dir(): raise ValueError("Not a directory")

    result = {
        "path": str(p), "total_size": 0, "total_files": 0, "total_dirs": 0,
        "by_extension": {}, "by_category": {}, "top_folders": [],
        "timeout": False, "analyzed_at": datetime.now().isoformat()
    }
    start_time = time.time()

    # Стек для итеративного DFS: (cur_dir_path, depth, top_folder_path)
    # top_folder_path хранит путь до папки верхнего уровня (глубина 1), к которой относится текущая ветка
    stack = []
    top_folder_stats = {}  # path -> stats dict

    # Инициализация: сканируем корень, добавляем в стек вложенные папки и файлы
    try:
        for entry in os.scandir(str(p)):
            if entry.name in ('System Volume Information', '$Recycle.Bin', 'Recovery'): continue
            if entry.is_dir(follow_symlinks=False):
                top_folder_stats[entry.path] = {"name": entry.name, "path": entry.path, "size": 0, "size_human": "", "files": 0, "dirs": 1}
                stack.append((entry.path, 1, entry.path))
                result["total_dirs"] += 1
            elif entry.is_file(follow_symlinks=False):
                try:
                    size = entry.stat(follow_symlinks=False).st_size
                    result["total_size"] += size
                    result["total_files"] += 1
                except: pass
    except (PermissionError, OSError): pass

    # Итеративный обход
    while stack:
        if time.time() - start_time > ANALYSIS_TIMEOUT:
            result["timeout"] = True
            break

        cur_dir, depth, top_folder_path = stack.pop()
        if depth > max_depth:
            continue

        try:
            entries = os.scandir(cur_dir)
        except (PermissionError, OSError):
            continue

        for entry in entries:
            if entry.name in ('System Volume Information', '$Recycle.Bin', 'Recovery'): continue

            if entry.is_dir(follow_symlinks=False):
                result["total_dirs"] += 1
                if top_folder_path and top_folder_path in top_folder_stats:
                    top_folder_stats[top_folder_path]["dirs"] += 1
                if depth + 1 <= max_depth:
                    stack.append((entry.path, depth + 1, top_folder_path))
            elif entry.is_file(follow_symlinks=False):
                try:
                    size = entry.stat(follow_symlinks=False).st_size
                    result["total_size"] += size
                    result["total_files"] += 1

                    # Обновляем статистику папки верхнего уровня
                    if top_folder_path and top_folder_path in top_folder_stats:
                        top_folder_stats[top_folder_path]["size"] += size
                        top_folder_stats[top_folder_path]["files"] += 1

                    # Расширения
                    ext = Path(entry.name).suffix.lower() or "no_ext"
                    if ext not in result["by_extension"]: result["by_extension"][ext] = {"count": 0, "size": 0}
                    result["by_extension"][ext]["count"] += 1
                    result["by_extension"][ext]["size"] += size
                    result["by_extension"][ext]["size_human"] = _format_size(result["by_extension"][ext]["size"])

                    # Категории
                    cat = _categorize(entry.path)
                    if cat not in result["by_category"]: result["by_category"][cat] = {"count": 0, "size": 0}
                    result["by_category"][cat]["count"] += 1
                    result["by_category"][cat]["size"] += size
                except:
                    pass

    # Формируем финальные данные
    result["total_size_human"] = _format_size(result["total_size"])
    result["top_folders"] = sorted(top_folder_stats.values(), key=lambda x: x["size"], reverse=True)[:10]
    for tf in result["top_folders"]:
        tf["size_human"] = _format_size(tf["size"])

    result["by_extension"] = dict(sorted(result["by_extension"].items(), key=lambda x: x[1]["size"], reverse=True))
    result["by_category"] = dict(sorted(result["by_category"].items(), key=lambda x: x[1]["size"], reverse=True))

    conversation_memory.add(op="analyze_directory", paths={"path": str(p)}, status="completed", dialog=dialog_ctx.get(), context=f"Analyzed {p}: {result['total_files']} files, {result['total_size_human']}")
    return result

# ─── File Info & Hash ────────────────────────────────────────────────────────
def get_file_info(path: str, include_hash=False):
    p = Path(normalize_path(path))
    try: _ensure_allowed(p, "get_file_info")
    except PermissionError as e: return {"error": str(e)}
    if not p.exists(): return {"path": str(p), "exists": False}
    stat = p.stat()
    res = {"path": str(p), "exists": True, "is_file": p.is_file(), "is_dir": p.is_dir(), "size": stat.st_size, "size_human": _format_size(stat.st_size), "created": datetime.fromtimestamp(stat.st_ctime).isoformat(), "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(), "accessed": datetime.fromtimestamp(stat.st_atime).isoformat()}
    if include_hash and p.is_file():
        algo = include_hash if isinstance(include_hash, str) else "md5"
        h = hashlib.md5() if algo == "md5" else hashlib.sha256()
        with open(p, 'rb') as f:
            for chunk in iter(lambda: f.read(65536), b''): h.update(chunk)
        res["hash"] = {"algorithm": algo, "value": h.hexdigest()}
    return res

def get_file_hash(path: str, algorithm="md5"):
    p = Path(normalize_path(path))
    try: _ensure_allowed(p, "get_file_hash")
    except PermissionError as e: return {"error": str(e)}
    if not p.is_file(): raise FileNotFoundError(f"File not found: {path}")
    h = hashlib.md5() if algorithm == "md5" else hashlib.sha256()
    with open(p, 'rb') as f:
        for chunk in iter(lambda: f.read(65536), b''): h.update(chunk)
    return {"path": str(p), "algorithm": algorithm, "hash": h.hexdigest()}

def compare_files(path1: str, path2: str, method="hash"):
    p1, p2 = Path(normalize_path(path1)), Path(normalize_path(path2))
    try: _ensure_allowed(p1, "compare_files"); _ensure_allowed(p2, "compare_files")
    except PermissionError as e: return {"error": str(e)}
    if method == "hash":
        h1 = get_file_hash(path1)["hash"]; h2 = get_file_hash(path2)["hash"]
        return {"identical": h1 == h2, "hash1": h1, "hash2": h2, "method": "md5"}
    size1 = p1.stat().st_size if p1.exists() else 0
    size2 = p2.stat().st_size if p2.exists() else 0
    if size1 != size2: return {"identical": False, "method": "content", "reason": "different sizes"}
    if size1 > MAX_READ_BYTES: return {"identical": False, "method": "content", "error": f"Files exceed {MAX_READ_BYTES} bytes limit", "size1": size1, "size2": size2}
    with open(p1, 'rb') as f1, open(p2, 'rb') as f2:
        while True:
            b1, b2 = f1.read(65536), f2.read(65536)
            if b1 != b2: return {"identical": False, "method": "content", "reason": "content mismatch"}
            if not b1: break
    return {"identical": True, "method": "content", "size1": size1, "size2": size2}

# ─── Server Setup ────────────────────────────────────────────────────────────
server = BaseMCPServer("filesystem-search", "4.2")
server.register_tool("search_files", {"description": "Search files by mask (multiple patterns via ;). Sync non-blocking.", "inputSchema": {"type": "object", "properties": {"path": {"type": "string"}, "pattern": {"type": "string"}, "recursive": {"type": "boolean", "default": True}, "max_files": {"type": "integer", "default": 10000}, "chunk_id": {"type": "string"}}, "required": ["path", "pattern"]}}, lambda **kw: search_files(kw["path"], kw["pattern"], kw.get("recursive", True), kw.get("max_files", 10000), kw.get("chunk_id")))
server.register_tool("search_content", {"description": "Search text inside files (simple mode)", "inputSchema": {"type": "object", "properties": {"path": {"type": "string"}, "query": {"type": "string"}, "extensions": {"type": "array", "items": {"type": "string"}}, "case_sensitive": {"type": "boolean", "default": False}, "max_files": {"type": "integer", "default": 500}}, "required": ["path", "query"]}}, lambda **kw: search_content(kw["path"], kw["query"], kw.get("extensions"), kw.get("case_sensitive", False), kw.get("max_files", 500)))
server.register_tool("search_content_advanced", {"description": "Advanced search with regex support and context lines", "inputSchema": {"type": "object", "properties": {"path": {"type": "string"}, "query": {"type": "string"}, "use_regex": {"type": "boolean", "default": False}, "case_sensitive": {"type": "boolean", "default": False}, "context_lines": {"type": "integer", "default": 0}, "max_files": {"type": "integer", "default": 100}, "extensions": {"type": "array", "items": {"type": "string"}}, "verbose": {"type": "boolean", "default": False}}, "required": ["path", "query"]}}, lambda **kw: search_content_advanced(kw["path"], kw["query"], kw.get("use_regex", False), kw.get("case_sensitive", False), kw.get("context_lines", 0), kw.get("max_files", 100), kw.get("extensions"), kw.get("verbose", False)))
server.register_tool("find_duplicates", {"description": "Find duplicates by hash or size", "inputSchema": {"type": "object", "properties": {"path": {"type": "string"}, "by": {"type": "string", "enum": ["hash", "size"], "default": "hash"}, "min_size_kb": {"type": "integer", "default": 1}, "max_size_mb": {"type": "integer", "default": 500}, "extensions": {"type": "array", "items": {"type": "string"}}}, "required": ["path"]}}, lambda **kw: find_duplicates(kw["path"], kw.get("by", "hash"), kw.get("min_size_kb", 1), kw.get("max_size_mb", 500), kw.get("extensions")))
server.register_tool("get_file_tree", {"description": "Directory tree (iterative, safe for deep nesting)", "inputSchema": {"type": "object", "properties": {"path": {"type": "string"}, "max_depth": {"type": "integer", "default": 5}, "include_files": {"type": "boolean", "default": False}}, "required": ["path"]}}, lambda **kw: get_file_tree(kw["path"], kw.get("max_depth", 5), kw.get("include_files", False)))
server.register_tool("analyze_directory", {"description": "Analyze directory: size, extensions, top folders", "inputSchema": {"type": "object", "properties": {"path": {"type": "string"}, "max_depth": {"type": "integer", "default": 3}, "group_by": {"type": "string", "enum": ["extension", "category", "date"], "default": "extension"}}, "required": ["path"]}}, lambda **kw: analyze_directory(kw["path"], kw.get("max_depth", 3), kw.get("group_by", "extension")))
server.register_tool("get_file_info", {"description": "File metadata", "inputSchema": {"type": "object", "properties": {"path": {"type": "string"}, "include_hash": {"type": ["boolean", "string"], "default": False}}, "required": ["path"]}}, lambda **kw: get_file_info(kw["path"], kw.get("include_hash", False)))
server.register_tool("get_file_hash", {"description": "Calculate file hash", "inputSchema": {"type": "object", "properties": {"path": {"type": "string"}, "algorithm": {"type": "string", "enum": ["md5", "sha256"], "default": "md5"}}, "required": ["path"]}}, lambda **kw: get_file_hash(kw["path"], kw.get("algorithm", "md5")))
server.register_tool("compare_files", {"description": "Compare two files", "inputSchema": {"type": "object", "properties": {"path1": {"type": "string"}, "path2": {"type": "string"}, "method": {"type": "string", "enum": ["hash", "content"], "default": "hash"}}, "required": ["path1", "path2"]}}, lambda **kw: compare_files(kw["path1"], kw["path2"], kw.get("method", "hash")))
server.register_tool("search_by_keywords", {
    "description": "Find files containing ALL specified keywords (AND logic)",
    "inputSchema": {
        "type": "object",
        "properties": {
            "path": {"type": "string"},
            "keywords": {"type": "array", "items": {"type": "string"}},
            "extensions": {"type": "array", "items": {"type": "string"}},
            "case_sensitive": {"type": "boolean", "default": False},
            "max_files": {"type": "integer", "default": 50}
        },
        "required": ["path", "keywords"]
    }
}, lambda **kw: search_by_keywords(
    kw["path"], kw["keywords"], kw.get("extensions"),
    kw.get("case_sensitive", False), kw.get("max_files", 50)
))

# ─── ДОПОЛНЕНИЕ: индексация содержимого офисных файлов ─────────────────────
import hashlib
import sqlite3
from mcp_rate_limiter import safe_call
from mcp_office_reader import DOCX_AVAILABLE, XLSX_AVAILABLE, PPTX_AVAILABLE
def index_office_files(path: str, extensions: List[str] = None, force_reindex: bool = False) -> Dict:
    """Индексирует содержимое DOCX/XLSX/PPTX для быстрого полнотекстового поиска."""
    service = "office_index"
    def _index():
        p = Path(normalize_path(path))
        _ensure_allowed(p, "index_office_files")
        if not p.is_dir():
            return {"error": "Path is not a directory"}
        index_db = p / ".office_index.db"
        conn = sqlite3.connect(str(index_db))
        conn.execute("""
        CREATE TABLE IF NOT EXISTS office_index (
            file_path TEXT PRIMARY KEY,
            content TEXT,
            hash TEXT,
            indexed_at REAL
        )
        """)
        conn.execute("CREATE VIRTUAL TABLE IF NOT EXISTS office_fts USING fts5(content, content=office_index)")
        exts = extensions or [".docx", ".xlsx", ".pptx"]
        indexed = 0
        skipped = 0
        for ext in exts:
            for f in p.rglob(f"*{ext}"):
                try:
                    current_hash = hashlib.md5(str(f).encode()).hexdigest()
                    row = conn.execute("SELECT hash FROM office_index WHERE file_path = ?", (str(f),)).fetchone()
                    if not force_reindex and row and row[0] == current_hash:
                        skipped += 1
                        continue
                    text = ""
                    if ext == ".docx" and DOCX_AVAILABLE:
                        from docx import Document
                        doc = Document(str(f))
                        text = " ".join(p.text for p in doc.paragraphs)
                    elif ext == ".xlsx" and XLSX_AVAILABLE:
                        from openpyxl import load_workbook
                        wb = load_workbook(str(f), read_only=True)
                        for sheet in wb.worksheets:
                            for row in sheet.iter_rows(values_only=True):
                                text += " ".join(str(cell) for cell in row if cell) + " "
                        wb.close()
                    elif ext == ".pptx" and PPTX_AVAILABLE:
                        from pptx import Presentation
                        prs = Presentation(str(f))
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text += shape.text + " "
                    if text:
                        conn.execute("INSERT OR REPLACE INTO office_index (file_path, content, hash, indexed_at) VALUES (?, ?, ?, ?)",
                                     (str(f), text[:50000], current_hash, time.time()))
                        indexed += 1
                except Exception as e:
                    _log(f"Index error {f}: {e}")
        conn.commit()
        conn.close()
        return {"indexed": indexed, "skipped": skipped, "db_path": str(index_db)}
    return safe_call(service, _index)

def search_indexed_office_files(path: str, query: str, limit: int = 50) -> Dict:
    """Полнотекстовый поиск по проиндексированным офисным файлам."""
    service = "office_search_index"
    def _search():
        p = Path(normalize_path(path))
        index_db = p / ".office_index.db"
        if not index_db.exists():
            return {"error": "Индекс не создан. Запустите index_office_files сначала."}
        conn = sqlite3.connect(str(index_db))
        try:
            cur = conn.execute("""
            SELECT file_path, content
            FROM office_fts
            JOIN office_index ON office_fts.rowid = office_index.rowid
            WHERE office_fts MATCH ?
            LIMIT ?
            """, (query, limit))
            results = [{"file": r[0], "snippet": r[1][:200]} for r in cur.fetchall()]
            return {"query": query, "results": results, "count": len(results)}
        finally:
            conn.close()
    return safe_call(service, _search)

server.register_tool("index_office_files", {
    "description": "Индексировать содержимое офисных файлов в папке",
    "inputSchema": {"type": "object", "properties": {"path": {"type": "string"}, "extensions": {"type": "array"}, "force_reindex": {"type": "boolean", "default": False}}, "required": ["path"]}
}, lambda **kw: index_office_files(kw["path"], kw.get("extensions"), kw.get("force_reindex", False)))
server.register_tool("search_office_index", {
    "description": "Полнотекстовый поиск по проиндексированным офисным файлам",
    "inputSchema": {"type": "object", "properties": {"path": {"type": "string"}, "query": {"type": "string"}, "limit": {"type": "integer", "default": 50}}, "required": ["path", "query"]}
}, lambda **kw: search_indexed_office_files(kw["path"], kw["query"], kw.get("limit", 50)))

if __name__ == "__main__":
    server.run()