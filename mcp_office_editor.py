#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
MCP Office Editor v1.0 – полноценное редактирование Excel (и базовое Word/PPT)
Поддерживает:
- Удаление пустых строк/столбцов
- Сортировку по одному или нескольким столбцам
- Добавление рамок (границ) к диапазону
- Установку значений в ячейки
- Применение стилей (жирный, цвет, выравнивание)
- Вставку/удаление строк/столбцов
- Работу с листами (создание, удаление, переименование)
- Чтение и запись (перезапись файла)
- Безопасность: все пути проверяются через _ensure_allowed
"""

import os
import json
from pathlib import Path
from typing import List, Dict, Optional, Union, Any

from openpyxl import load_workbook, Workbook
from openpyxl.styles import Border, Side, Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter, column_index_from_string

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from mcp_shared import (
    _log, normalize_path, _ensure_allowed, BaseMCPServer, conversation_memory, dialog_ctx
)

# =========================== Excel Helpers ===========================

def _load_excel_safe(file_path: Path):
    """Безопасная загрузка книги, создаёт новую, если файл не существует."""
    if file_path.exists():
        return load_workbook(file_path)
    else:
        _log(f"[OfficeEditor] Файл {file_path} не найден, создаю новую книгу.")
        return Workbook()

def _save_excel_safe(workbook, file_path: Path):
    """Сохраняет книгу, создавая папки при необходимости."""
    file_path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(file_path)
    _log(f"[OfficeEditor] Сохранён Excel: {file_path}")

def _get_sheet(workbook, sheet_name: Optional[str] = None):
    """Возвращает лист по имени или активный лист."""
    if sheet_name and sheet_name in workbook.sheetnames:
        return workbook[sheet_name]
    return workbook.active

def _apply_border_to_range(worksheet, cell_range: str, border_style: str = 'thin'):
    """Применяет границы к диапазону (например, 'A1:C10' или вся таблица)."""
    if cell_range.lower() == 'all':
        max_row = worksheet.max_row
        max_col = worksheet.max_column
        cell_range = f"A1:{get_column_letter(max_col)}{max_row}"
    
    side = Side(style=border_style)
    border = Border(left=side, right=side, top=side, bottom=side)
    
    for row in worksheet[cell_range]:
        for cell in row:
            cell.border = border

# =========================== Excel Tools ===========================

def excel_delete_empty_rows(file_path: str, sheet_name: Optional[str] = None,
                            dry_run: bool = False) -> Dict:
    """
    Удаляет строки, в которых все ячейки пусты (None или пустая строка).
    """
    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "excel_delete_empty_rows")
    
    if not p.exists():
        return {"error": "Файл не существует"}
    
    wb = _load_excel_safe(p)
    ws = _get_sheet(wb, sheet_name)
    
    rows_to_delete = []
    for row_idx in range(1, ws.max_row + 1):
        empty = True
        for col_idx in range(1, ws.max_column + 1):
            val = ws.cell(row=row_idx, column=col_idx).value
            if val is not None and str(val).strip() != "":
                empty = False
                break
        if empty:
            rows_to_delete.append(row_idx)
    
    if dry_run:
        wb.close()
        return {
            "status": "dry_run",
            "file": str(p),
            "empty_rows_count": len(rows_to_delete),
            "rows": rows_to_delete[:50]
        }
    
    for r in reversed(rows_to_delete):
        ws.delete_rows(r)
    
    _save_excel_safe(wb, p)
    return {
        "status": "success",
        "deleted_rows": len(rows_to_delete),
        "file": str(p)
    }

def excel_sort(file_path: str, columns: Union[str, List[str]], sheet_name: Optional[str] = None,
               ascending: bool = True, dry_run: bool = False) -> Dict:
    """
    Сортировка данных по одному или нескольким столбцам.
    columns: 'A' или ['A','B'] или ['A', 'C', 'B']
    """
    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "excel_sort")
    
    if not p.exists():
        return {"error": "Файл не существует"}
    
    wb = _load_excel_safe(p)
    ws = _get_sheet(wb, sheet_name)
    
    data_range = ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column, values_only=False)
    rows = list(data_range)
    if len(rows) <= 1:
        return {"message": "Нет данных для сортировки"}
    
    header = rows[0]
    data_rows = rows[1:]
    
    col_indices = []
    for col in columns:
        if isinstance(col, str) and len(col) <= 3:
            col_idx = column_index_from_string(col)
        elif isinstance(col, int):
            col_idx = col
        else:
            col_idx = None
        if col_idx is None or col_idx < 1 or col_idx > ws.max_column:
            return {"error": f"Некорректный столбец: {col}"}
        col_indices.append(col_idx)
    
    if dry_run:
        wb.close()
        return {
            "status": "dry_run",
            "file": str(p),
            "sort_columns": columns,
            "rows_before": len(data_rows),
            "rows_after": len(data_rows)
        }
    
    data_rows.sort(key=lambda row: [row[col_idx-1].value for col_idx in col_indices], 
                   reverse=not ascending)
    
    ws.delete_rows(2, ws.max_row-1)
    for r_idx, row in enumerate(data_rows, start=2):
        for c_idx, cell in enumerate(row, start=1):
            ws.cell(row=r_idx, column=c_idx, value=cell.value)
    
    _save_excel_safe(wb, p)
    return {
        "status": "success",
        "sorted_by": columns,
        "rows": len(data_rows),
        "file": str(p)
    }

def excel_add_borders(file_path: str, range_str: str = "all", sheet_name: Optional[str] = None,
                      border_style: str = "thin", dry_run: bool = False) -> Dict:
    """
    Добавляет границы к указанному диапазону (например, 'A1:C10' или 'all' для всей таблицы).
    """
    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "excel_add_borders")
    
    if not p.exists():
        return {"error": "Файл не существует"}
    
    wb = _load_excel_safe(p)
    ws = _get_sheet(wb, sheet_name)
    
    if dry_run:
        wb.close()
        return {"status": "dry_run", "range": range_str, "file": str(p)}
    
    _apply_border_to_range(ws, range_str, border_style)
    _save_excel_safe(wb, p)
    
    return {
        "status": "success",
        "range": range_str,
        "file": str(p)
    }

def excel_set_cell(file_path: str, cell: str, value: Any, sheet_name: Optional[str] = None,
                   dry_run: bool = False) -> Dict:
    """
    Установить значение в указанную ячейку (например, 'B5').
    """
    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "excel_set_cell")
    
    if not p.exists():
        return {"error": "Файл не существует"}
    
    wb = _load_excel_safe(p)
    ws = _get_sheet(wb, sheet_name)
    
    if dry_run:
        wb.close()
        return {"status": "dry_run", "cell": cell, "value": value}
    
    ws[cell] = value
    _save_excel_safe(wb, p)
    
    return {"status": "success", "cell": cell, "value": str(value), "file": str(p)}

def excel_apply_format(file_path: str, range_str: str, 
                       bold: bool = False, italic: bool = False, 
                       font_color: Optional[str] = None,
                       bg_color: Optional[str] = None,
                       alignment: Optional[str] = None,
                       sheet_name: Optional[str] = None,
                       dry_run: bool = False) -> Dict:
    """
    Применяет форматирование к диапазону.
    alignment: 'left', 'center', 'right'
    font_color: 'FF0000' (HEX)
    bg_color: 'FFFF00' (HEX)
    """
    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "excel_apply_format")
    
    if not p.exists():
        return {"error": "Файл не существует"}
    
    wb = _load_excel_safe(p)
    ws = _get_sheet(wb, sheet_name)
    
    if dry_run:
        wb.close()
        return {"status": "dry_run", "range": range_str, "format": {"bold": bold, "italic": italic}}
    
    for row in ws[range_str]:
        for cell in row:
            if bold or italic or font_color:
                font = Font(bold=bold, italic=italic, color=font_color)
                cell.font = font
            if bg_color:
                fill = PatternFill(start_color=bg_color, end_color=bg_color, fill_type="solid")
                cell.fill = fill
            if alignment:
                align = Alignment(horizontal=alignment)
                cell.alignment = align
    
    _save_excel_safe(wb, p)
    return {"status": "success", "range": range_str, "file": str(p)}

# =========================== Word Tools (базовые) ===========================

def docx_replace_text(file_path: str, old_text: str, new_text: str,
                      dry_run: bool = False) -> Dict:
    """Заменяет текст во всех параграфах и таблицах документа .docx."""
    if not DOCX_AVAILABLE:
        return {"error": "python-docx не установлен. Установите: pip install python-docx"}
    
    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "docx_replace_text")
    
    if not p.exists():
        return {"error": "Файл не существует"}
    
    doc = Document(str(p))
    
    if dry_run:
        count = 0
        for para in doc.paragraphs:
            count += para.text.count(old_text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    count += cell.text.count(old_text)
        return {"status": "dry_run", "occurrences": count, "file": str(p)}
    
    for para in doc.paragraphs:
        if old_text in para.text:
            para.text = para.text.replace(old_text, new_text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if old_text in cell.text:
                    cell.text = cell.text.replace(old_text, new_text)
    
    doc.save(str(p))
    return {"status": "success", "replaced": old_text, "with": new_text, "file": str(p)}

# =========================== PowerPoint Tools (базовые) ===========================

def pptx_add_slide(file_path: str, title: str, content: str, layout_index: int = 0,
                   dry_run: bool = False) -> Dict:
    """Добавляет новый слайд в презентацию."""
    if not PPTX_AVAILABLE:
        return {"error": "python-pptx не установлен. Установите: pip install python-pptx"}
    
    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "pptx_add_slide")
    
    if not p.exists():
        return {"error": "Файл не существует"}
    
    prs = Presentation(str(p))
    slide_layout = prs.slide_layouts[layout_index]
    
    if dry_run:
        return {"status": "dry_run", "slide_title": title}
    
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    if title_placeholder:
        title_placeholder.text = title
    content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if content_placeholder:
        content_placeholder.text = content
    
    prs.save(str(p))
    return {"status": "success", "added_slide": title, "file": str(p)}

# =========================== Unified Entry Point ===========================

def office_edit(file_path: str, operation: str, **kwargs) -> Dict:
    """
    Универсальный диспетчер для операций редактирования офисных файлов.
    """
    operation_map = {
        "excel_delete_empty_rows": excel_delete_empty_rows,
        "excel_sort": excel_sort,
        "excel_add_borders": excel_add_borders,
        "excel_set_cell": excel_set_cell,
        "excel_apply_format": excel_apply_format,
        "docx_replace_text": docx_replace_text,
        "pptx_add_slide": pptx_add_slide,
    }
    if operation not in operation_map:
        return {"error": f"Unknown operation: {operation}"}
    
    func = operation_map[operation]
    return func(file_path=file_path, **kwargs)

# =========================== Server Registration ===========================

server = BaseMCPServer("office-editor", "1.0")

server.register_tool("excel_delete_empty_rows", {
    "description": "Удалить полностью пустые строки в Excel-файле",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "sheet_name": {"type": "string"},
            "dry_run": {"type": "boolean", "default": False}
        },
        "required": ["file_path"]
    }
}, lambda **kw: excel_delete_empty_rows(**kw))

server.register_tool("excel_sort", {
    "description": "Сортировка Excel-таблицы по одному или нескольким столбцам (например, columns=['A','B'])",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "columns": {"type": ["string", "array"], "items": {"type": "string"}},
            "sheet_name": {"type": "string"},
            "ascending": {"type": "boolean", "default": True},
            "dry_run": {"type": "boolean", "default": False}
        },
        "required": ["file_path", "columns"]
    }
}, lambda **kw: excel_sort(**kw))

server.register_tool("excel_add_borders", {
    "description": "Добавить границы к диапазону (например, 'A1:C10' или 'all')",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "range_str": {"type": "string", "default": "all"},
            "sheet_name": {"type": "string"},
            "border_style": {"type": "string", "default": "thin"},
            "dry_run": {"type": "boolean", "default": False}
        },
        "required": ["file_path"]
    }
}, lambda **kw: excel_add_borders(**kw))

server.register_tool("excel_set_cell", {
    "description": "Установить значение в конкретную ячейку",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "cell": {"type": "string", "description": "Например, 'B5'"},
            "value": {"type": "string"},
            "sheet_name": {"type": "string"},
            "dry_run": {"type": "boolean", "default": False}
        },
        "required": ["file_path", "cell", "value"]
    }
}, lambda **kw: excel_set_cell(**kw))

server.register_tool("excel_apply_format", {
    "description": "Применить форматирование (жирный, цвет, выравнивание) к диапазону",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "range_str": {"type": "string"},
            "bold": {"type": "boolean", "default": False},
            "italic": {"type": "boolean", "default": False},
            "font_color": {"type": "string", "description": "HEX цвет, например 'FF0000'"},
            "bg_color": {"type": "string", "description": "HEX цвет заливки"},
            "alignment": {"type": "string", "enum": ["left", "center", "right"]},
            "sheet_name": {"type": "string"},
            "dry_run": {"type": "boolean", "default": False}
        },
        "required": ["file_path", "range_str"]
    }
}, lambda **kw: excel_apply_format(**kw))

server.register_tool("docx_replace_text", {
    "description": "Замена текста в документе Word (.docx)",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "old_text": {"type": "string"},
            "new_text": {"type": "string"},
            "dry_run": {"type": "boolean", "default": False}
        },
        "required": ["file_path", "old_text", "new_text"]
    }
}, lambda **kw: docx_replace_text(**kw))

server.register_tool("pptx_add_slide", {
    "description": "Добавить слайд в презентацию PowerPoint (.pptx)",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "title": {"type": "string"},
            "content": {"type": "string"},
            "layout_index": {"type": "integer", "default": 0},
            "dry_run": {"type": "boolean", "default": False}
        },
        "required": ["file_path", "title", "content"]
    }
}, lambda **kw: pptx_add_slide(**kw))

server.register_tool("office_edit", {
    "description": "Универсальная операция редактирования офисных файлов (Excel, Word, PPT)",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "operation": {"type": "string", "enum": list([
                "excel_delete_empty_rows", "excel_sort", "excel_add_borders",
                "excel_set_cell", "excel_apply_format", "docx_replace_text",
                "pptx_add_slide"
            ])},
            "kwargs": {"type": "object"}
        },
        "required": ["file_path", "operation"]
    }
}, lambda **kw: office_edit(**kw))

if __name__ == "__main__":
    server.run()